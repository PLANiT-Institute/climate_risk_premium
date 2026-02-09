#!/usr/bin/env python3
"""Generate CRP PowerPoint presentation.

Creates an 11-slide presentation covering:
  1. Cashflow assumptions (plant, revenue, costs)
  2. Physical risk methodology (PLANiT framework, map, hazard logic)
  3. Key results (scenario comparison, rating migration, findings)

Usage:
    python3 scripts/generate_presentation.py

Output:
    results/CRP_Presentation.pptx
    results/figures/fig_samcheok_map.png  (generated map)
"""

import csv
import json
import os
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyArrowPatch
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ────────────────────────────────────────────────────────────────
ROOT = Path(__file__).resolve().parent.parent
PLANT_CSV = ROOT / "data" / "raw" / "plant.csv"
SCENARIO_CSV = ROOT / "data" / "processed" / "scenario_comparison.csv"
GEOJSON = ROOT / "crp-dashboard" / "public" / "data" / "geo" / "samcheok_power_grid.geojson"
FIG1 = ROOT / "results" / "figures" / "fig1_npv_comparison.png"
FIG3 = ROOT / "results" / "figures" / "fig3_rating_migration.png"
MAP_OUT = ROOT / "results" / "figures" / "fig_samcheok_map.png"
OUTPUT_PPTX = ROOT / "results" / "CRP_Presentation.pptx"

# ── Colors ───────────────────────────────────────────────────────────────
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x00, 0x7B, 0x7F)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREEN_ACCENT = RGBColor(0x27, 0xAE, 0x60)


# ═════════════════════════════════════════════════════════════════════════
# Helper: load CSV as list of dicts
# ═════════════════════════════════════════════════════════════════════════
def load_csv(path):
    with open(path, newline="", encoding="utf-8") as f:
        return list(csv.DictReader(f))


def load_plant_params(path):
    """Return dict {parameter: value}."""
    rows = load_csv(path)
    return {r["parameter"]: r["value"] for r in rows}


# ═════════════════════════════════════════════════════════════════════════
# Map generation
# ═════════════════════════════════════════════════════════════════════════
def generate_map():
    """Create a schematic Samcheok site map from GeoJSON."""
    with open(GEOJSON, encoding="utf-8") as f:
        geo = json.load(f)

    # Extract plant polygon
    plant_feature = None
    substation = None
    transmission_line = None
    plant_route = None

    for feat in geo["features"]:
        props = feat["properties"]
        if feat["geometry"]["type"] == "Polygon":
            plant_feature = feat
        elif feat["geometry"]["type"] == "Point":
            substation = feat
        elif feat["geometry"]["type"] == "LineString":
            if "voltage_kv" in props:
                transmission_line = feat
            else:
                plant_route = feat

    # Plant polygon coords
    poly_coords = plant_feature["geometry"]["coordinates"][0]
    plant_lons = [c[0] for c in poly_coords]
    plant_lats = [c[1] for c in poly_coords]
    center_lon = plant_feature["properties"]["center_lon"]
    center_lat = plant_feature["properties"]["center_lat"]

    fig, ax = plt.subplots(1, 1, figsize=(10, 8))

    # Background - light blue for sea (east of plant)
    ax.set_facecolor("#E8F4F8")

    # Draw simplified coastline context (schematic land area to west)
    # Wider land mass covering the left portion of the view
    land_x = [129.130, 129.130, 129.185, 129.195, 129.200, 129.200,
              129.195, 129.185, 129.175, 129.170, 129.168, 129.165,
              129.162, 129.158, 129.155, 129.150, 129.145, 129.140, 129.130]
    land_y = [37.365, 37.450, 37.450, 37.440, 37.420, 37.415,
              37.400, 37.398, 37.400, 37.398, 37.396, 37.397,
              37.396, 37.395, 37.392, 37.388, 37.380, 37.370, 37.365]
    ax.fill(land_x, land_y, color="#D5E8D4", alpha=0.6, zorder=1)

    # Plant polygon (red fill)
    ax.fill(plant_lons, plant_lats, color="#C0392B", alpha=0.35, zorder=3)
    ax.plot(plant_lons, plant_lats, color="#C0392B", linewidth=2, zorder=4)

    # Risk zone circle (25 km radius ≈ 0.225° at this latitude)
    risk_circle = plt.Circle(
        (center_lon, center_lat), 0.03, fill=False,
        color="#E67E22", linewidth=2, linestyle="--", zorder=5, label="Risk zone (3 km)"
    )
    ax.add_patch(risk_circle)

    # Plant center marker
    ax.plot(center_lon, center_lat, "^", color="#C0392B", markersize=14, zorder=6)
    ax.annotate(
        "Samcheok Coal Plant\n2,100 MW",
        (center_lon, center_lat),
        textcoords="offset points", xytext=(15, 15),
        fontsize=11, fontweight="bold", color="#1B2A4A",
        bbox=dict(boxstyle="round,pad=0.3", facecolor="white", edgecolor="#1B2A4A", alpha=0.9),
        arrowprops=dict(arrowstyle="->", color="#1B2A4A"),
        zorder=7
    )

    # Transmission route (if available, show first few segments only)
    if plant_route:
        route_coords = plant_route["geometry"]["coordinates"]
        # Show first 5 segments for context
        route_lons = [c[0] for c in route_coords[:5]]
        route_lats = [c[1] for c in route_coords[:5]]
        ax.plot(route_lons, route_lats, color="#7F8C8D", linewidth=1.5,
                linestyle="-.", alpha=0.6, zorder=2)

    # Sea label
    ax.text(129.205, 37.435, "East Sea\n(Sea of Japan)",
            fontsize=10, fontstyle="italic", color="#2980B9", ha="center", zorder=7)

    # Coordinate labels
    ax.set_xlabel("Longitude (°E)", fontsize=11)
    ax.set_ylabel("Latitude (°N)", fontsize=11)
    ax.set_title("Samcheok Coal-Fired Power Plant Site", fontsize=14, fontweight="bold",
                 color="#1B2A4A", pad=15)

    # Set extent around plant with buffer for risk zone
    ax.set_xlim(129.135, 129.215)
    ax.set_ylim(37.370, 37.445)

    # Force full coordinate labels (no offset notation)
    ax.xaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"{x:.3f}"))
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda y, _: f"{y:.3f}"))

    # Grid
    ax.grid(True, alpha=0.3, linestyle="--")
    ax.tick_params(labelsize=9)

    # North arrow
    ax.annotate("N", xy=(0.95, 0.95), xycoords="axes fraction",
                fontsize=14, fontweight="bold", ha="center", va="center",
                bbox=dict(boxstyle="round", facecolor="white", edgecolor="black"))
    ax.annotate("", xy=(0.95, 0.93), xycoords="axes fraction",
                xytext=(0.95, 0.87), textcoords="axes fraction",
                arrowprops=dict(arrowstyle="->", lw=2))

    # Scale bar (approx 1 km)
    # At 37.4°N, 1° lon ≈ 88 km, so 1km ≈ 0.01136°
    scale_lon_start = 129.140
    scale_lon_end = scale_lon_start + 0.01136
    scale_lat = 37.375
    ax.plot([scale_lon_start, scale_lon_end], [scale_lat, scale_lat],
            "k-", linewidth=3, zorder=8)
    ax.text((scale_lon_start + scale_lon_end) / 2, scale_lat - 0.001,
            "1 km", ha="center", fontsize=9, zorder=8)

    # Coordinate annotation
    ax.text(0.02, 0.02,
            f"Center: {center_lon:.3f}°E, {center_lat:.3f}°N",
            transform=ax.transAxes, fontsize=8, color="#666",
            bbox=dict(facecolor="white", alpha=0.8, edgecolor="none"))

    # Legend
    legend_elements = [
        mpatches.Patch(facecolor="#C0392B", alpha=0.35, edgecolor="#C0392B",
                       label="Plant footprint"),
        plt.Line2D([0], [0], color="#E67E22", linewidth=2, linestyle="--",
                   label="Risk zone (3 km)"),
        mpatches.Patch(facecolor="#D5E8D4", alpha=0.7, label="Land"),
    ]
    ax.legend(handles=legend_elements, loc="lower right", fontsize=9,
              framealpha=0.9)

    plt.tight_layout()
    MAP_OUT.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(MAP_OUT, dpi=200, bbox_inches="tight")
    plt.close(fig)
    print(f"  Map saved: {MAP_OUT}")


# ═════════════════════════════════════════════════════════════════════════
# Slide helpers
# ═════════════════════════════════════════════════════════════════════════
def set_slide_bg(slide, color=WHITE):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, text, left=0, top=0,
                  width=Inches(13.333), height=Inches(1.0)):
    """Add a colored title bar at the top of a slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=14, bold=False, color=DARK_GRAY,
                 alignment=PP_ALIGN.LEFT):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, items, left, top, width, height,
                    font_size=14, color=DARK_GRAY, spacing=Pt(6)):
    """Add bulleted text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = spacing
        p.level = 0
        # Bullet character
        p.text = f"\u2022  {item}"
    return txBox


def add_table(slide, rows, col_widths, left, top, font_size=12):
    """Add a styled table. rows[0] = header row."""
    n_rows = len(rows)
    n_cols = len(rows[0])
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, left, top,
        Inches(total_w), Inches(0.35 * n_rows)
    )
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)

    for r_idx, row_data in enumerate(rows):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(font_size)
            p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT

            if r_idx == 0:
                # Header row styling
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
                p.font.color.rgb = WHITE
                p.font.bold = True
            else:
                # Alternate row shading
                if r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GRAY
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE
                p.font.color.rgb = DARK_GRAY

    return table_shape


def add_accent_line(slide, top):
    """Add a thin teal accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), top, Inches(12.3), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()


# ═════════════════════════════════════════════════════════════════════════
# Slide builders
# ═════════════════════════════════════════════════════════════════════════
def make_title_slide(prs):
    """Slide 1: Title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, WHITE)

    # Large navy rectangle (upper half)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(4.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    # Title text
    add_text_box(slide, "Climate Risk Premium", Inches(1), Inches(1.2),
                 Inches(11), Inches(1), font_size=42, bold=True, color=WHITE)

    # Subtitle
    add_text_box(slide, "Samcheok Coal-Fired Power Plant",
                 Inches(1), Inches(2.2), Inches(11), Inches(0.7),
                 font_size=28, color=RGBColor(0xA0, 0xD0, 0xD0))

    # Accent line
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.1), Inches(3), Pt(4)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = TEAL
    shape2.line.fill.background()

    # Description
    add_text_box(slide,
                 "Cashflow Assumptions, Physical Risk Methodology & Key Results",
                 Inches(1), Inches(3.4), Inches(10), Inches(0.6),
                 font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC))

    # Date
    add_text_box(slide, "February 2026",
                 Inches(1), Inches(5.0), Inches(4), Inches(0.5),
                 font_size=16, color=DARK_GRAY)

    # Methodology note
    add_text_box(slide,
                 "Model: DCF + Monte Carlo  |  Physical Risk: CLIMADA + OS-Climate (PLANiT)  |  11 Scenarios",
                 Inches(1), Inches(5.6), Inches(11), Inches(0.5),
                 font_size=12, color=RGBColor(0x88, 0x88, 0x88))


def make_plant_overview(prs, plant):
    """Slide 2: Plant Overview & Key Parameters."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Plant Overview & Key Parameters")

    # Section label
    add_text_box(slide, "SECTION 1: CASHFLOW ASSUMPTIONS",
                 Inches(0.5), Inches(1.15), Inches(5), Inches(0.4),
                 font_size=11, bold=True, color=TEAL)

    # Plant parameters table
    table_data = [
        ["Parameter", "Value", "Unit"],
        ["Installed Capacity", "2,100", "MW"],
        ["Baseline Capacity Factor", "85%", "ratio"],
        ["Operating Life", "40", "years"],
        ["Total CAPEX", "4,900", "million USD"],
        ["Commercial Operation Date", "2024", "year"],
        ["Heat Rate", "9.5", "MMBtu/MWh"],
        ["CO\u2082 Emissions Intensity", "0.95", "tCO\u2082/MWh"],
        ["Base Outage Rate", "3%", ""],
    ]
    add_table(slide, table_data, [4.0, 2.0, 2.0], Inches(0.5), Inches(1.7), font_size=13)

    # Key callout box on the right
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(1.7),
        Inches(3.8), Inches(4.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
    shape.line.color.rgb = TEAL

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)

    p = tf.paragraphs[0]
    p.text = "Key Highlights"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY

    highlights = [
        "Largest coal plant in Korea",
        "Annual generation: ~15.6 TWh",
        "Baseline revenue: ~USD 1,252M/yr",
        "40-year DCF horizon",
        "Korea-specific climate parameters",
    ]
    for h in highlights:
        p = tf.add_paragraph()
        p.text = f"\u2022  {h}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)


def make_price_revenue(prs):
    """Slide 3: Electricity Price & Revenue."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Electricity Price & Revenue")

    add_text_box(slide, "SECTION 1: CASHFLOW ASSUMPTIONS",
                 Inches(0.5), Inches(1.15), Inches(5), Inches(0.4),
                 font_size=11, bold=True, color=TEAL)

    # Left column - Price assumptions
    add_text_box(slide, "Price Assumptions", Inches(0.5), Inches(1.7),
                 Inches(5.5), Inches(0.5), font_size=18, bold=True, color=NAVY)

    items = [
        "Base power price: 80 USD/MWh",
        "SMP trajectory: 85 \u2192 58 USD/MWh (2024\u20132050)",
        "Carbon price: 15 \u2192 150 USD/tCO\u2082 (2024\u20132050)",
        "Carbon cost = 0.95 tCO\u2082/MWh \u00d7 carbon price",
    ]
    add_bullet_list(slide, items, Inches(0.5), Inches(2.3),
                    Inches(5.5), Inches(3.0), font_size=13)

    # Right column - Revenue calculation
    add_text_box(slide, "Revenue Calculation", Inches(7.0), Inches(1.7),
                 Inches(5.5), Inches(0.5), font_size=18, bold=True, color=NAVY)

    rev_items = [
        "Annual generation = 2,100 MW \u00d7 0.85 CF \u00d7 8,760 h",
        "                         = 15,638,400 MWh/yr (~15.6 TWh)",
        "Baseline revenue = 15.6 TWh \u00d7 80 USD/MWh",
        "                           = ~USD 1,252 M/yr",
        "Revenue formula: Annual MWh \u00d7 (SMP \u2212 Carbon Cost)",
    ]
    add_bullet_list(slide, rev_items, Inches(7.0), Inches(2.3),
                    Inches(5.8), Inches(3.0), font_size=13)

    # Bottom callout
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.5),
        Inches(12.3), Inches(1.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFD, 0xF2, 0xE9)
    shape.line.color.rgb = RGBColor(0xE6, 0x7E, 0x22)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.15)

    p = tf.paragraphs[0]
    p.text = "Key Insight: Rising carbon costs progressively erode coal plant margins"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xE6, 0x7E, 0x22)

    p2 = tf.add_paragraph()
    p2.text = ("Under aggressive transition, carbon costs of 150 USD/tCO\u2082 "
               "add ~142 USD/MWh, far exceeding the SMP of 58 USD/MWh by 2050.")
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_GRAY


def make_cost_financing(prs):
    """Slide 4: Cost Structure & Financing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Cost Structure & Financing")

    add_text_box(slide, "SECTION 1: CASHFLOW ASSUMPTIONS",
                 Inches(0.5), Inches(1.15), Inches(5), Inches(0.4),
                 font_size=11, bold=True, color=TEAL)

    # Cost table
    cost_data = [
        ["Cost Component", "Value", "Annual Amount"],
        ["Fuel Cost", "9.5 MMBtu/MWh \u00d7 3.2 USD/MMBtu = 30.4 USD/MWh", "USD 475 M/yr"],
        ["Fixed O&M", "35 USD/kW-yr", "USD 73.5 M/yr"],
        ["Variable O&M", "4.5 USD/MWh", "USD 70.4 M/yr"],
        ["Total Operating Cost", "", "USD 619 M/yr"],
    ]
    add_table(slide, cost_data, [3.0, 5.0, 2.5], Inches(0.5), Inches(1.7), font_size=12)

    # Financing table
    add_text_box(slide, "Financing Structure", Inches(0.5), Inches(4.0),
                 Inches(5), Inches(0.5), font_size=18, bold=True, color=NAVY)

    fin_data = [
        ["Parameter", "Value"],
        ["Debt / Equity", "70% / 30%"],
        ["Debt Interest Rate", "6.1%"],
        ["Debt Tenor", "20 years"],
        ["Discount Rate (WACC)", "8.0%"],
        ["Tax Rate", "24%"],
        ["CAPEX", "USD 4,900 M"],
    ]
    add_table(slide, fin_data, [3.5, 3.0], Inches(0.5), Inches(4.6), font_size=12)

    # Right side - EBITDA callout
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.7),
        Inches(4.3), Inches(5.0)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
    shape.line.color.rgb = TEAL

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)

    p = tf.paragraphs[0]
    p.text = "Baseline Economics"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY

    econ = [
        ("Revenue", "USD 1,252 M/yr"),
        ("Operating Cost", "USD 619 M/yr"),
        ("EBITDA", "USD 633 M/yr"),
        ("EBITDA Margin", "~51%"),
        ("", ""),
        ("NPV (8% disc.)", "USD 3,103 M"),
        ("Project IRR", "12.0%"),
        ("Avg DSCR", "2.15x"),
    ]
    for label, val in econ:
        p = tf.add_paragraph()
        if label:
            p.text = f"{label}: {val}"
        else:
            p.text = ""
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)


def make_planit_framework(prs):
    """Slide 5: PLANiT Framework Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Physical Risk Framework: PLANiT Pipeline")

    add_text_box(slide, "SECTION 2: PHYSICAL RISK METHODOLOGY",
                 Inches(0.5), Inches(1.15), Inches(5), Inches(0.4),
                 font_size=11, bold=True, color=TEAL)

    # Pipeline diagram using shapes
    # Row 1: Data sources
    sources = [
        ("CLIMADA\n(Wildfire AAI)", Inches(0.5)),
        ("OS-Climate PhysRisk\n(Drought)", Inches(0.5)),
        ("OS-Climate PhysRisk\n(Water Risk)", Inches(0.5)),
    ]
    y_start = Inches(2.0)
    for i, (label, _) in enumerate(sources):
        x = Inches(0.5)
        y = y_start + Inches(i * 1.4)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.8), Inches(1.0)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
        shape.line.color.rgb = NAVY
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Arrow
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, Inches(3.5), y + Inches(0.25),
            Inches(1.0), Inches(0.5)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = TEAL
        arrow.line.fill.background()

    # PLANiT Adapter box
    adapter_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(2.3),
        Inches(2.5), Inches(3.0)
    )
    adapter_shape.fill.solid()
    adapter_shape.fill.fore_color.rgb = TEAL
    adapter_shape.line.color.rgb = NAVY
    tf = adapter_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PLANiT\nAdapter"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p2 = tf.add_paragraph()
    p2.text = "\nConverts raw\nhazard outputs to\ncashflow impacts"
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(0xE0, 0xF0, 0xF0)
    p2.alignment = PP_ALIGN.CENTER

    # Arrow from adapter
    arrow2 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(7.5), Inches(3.3),
        Inches(1.0), Inches(0.5)
    )
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = TEAL
    arrow2.line.fill.background()

    # PhysicalAdjustments box
    pa_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(2.3),
        Inches(2.5), Inches(3.0)
    )
    pa_shape.fill.solid()
    pa_shape.fill.fore_color.rgb = NAVY
    pa_shape.line.color.rgb = TEAL
    tf = pa_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Physical\nAdjustments"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p2 = tf.add_paragraph()
    p2.text = "\n\u2022 outage_rate\n\u2022 capacity_derate\n\u2022 capex_increase\n\u2022 insurance_increase"
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(0xA0, 0xD0, 0xD0)
    p2.alignment = PP_ALIGN.CENTER

    # Arrow to CF Impact
    arrow3 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(11.5), Inches(3.3),
        Inches(0.7), Inches(0.5)
    )
    arrow3.fill.solid()
    arrow3.fill.fore_color.rgb = RED_ACCENT
    arrow3.line.fill.background()

    # CF Impact
    cf_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(12.0), Inches(2.8),
        Inches(1.2), Inches(1.5)
    )
    cf_shape.fill.solid()
    cf_shape.fill.fore_color.rgb = RED_ACCENT
    cf_shape.line.fill.background()
    tf = cf_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "CF\nImpact"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Bottom note
    add_text_box(slide,
                 "Source: src/planit/adapter.py, src/planit/config.py  |  "
                 "SSP Mapping: baseline\u2192historical, moderate\u2192SSP1-2.6, high\u2192SSP5-8.5",
                 Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
                 font_size=10, color=RGBColor(0x88, 0x88, 0x88))


def make_map_slide(prs):
    """Slide 6: Samcheok Site Location Map."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Samcheok Site Location")

    add_text_box(slide, "SECTION 2: PHYSICAL RISK METHODOLOGY",
                 Inches(0.5), Inches(1.15), Inches(5), Inches(0.4),
                 font_size=11, bold=True, color=TEAL)

    # Embed map image
    if MAP_OUT.exists():
        slide.shapes.add_picture(
            str(MAP_OUT), Inches(0.8), Inches(1.7), Inches(7.5), Inches(5.5)
        )
    else:
        add_text_box(slide, "[Map image not found]",
                     Inches(2), Inches(3), Inches(6), Inches(1),
                     font_size=20, color=RED_ACCENT, alignment=PP_ALIGN.CENTER)

    # Info panel on right
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.7),
        Inches(4.0), Inches(5.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
    shape.line.color.rgb = NAVY

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)

    p = tf.paragraphs[0]
    p.text = "Site Details"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY

    details = [
        "Location: Samcheok, Gangwon Province",
        "Coordinates: 129.174\u00b0E, 37.409\u00b0N",
        "Coastal site on East Sea",
        "Connected via 765kV transmission",
        "Grid: Samcheok \u2192 Sintaebaek Substation",
        "",
        "Physical Risk Exposure:",
        "  \u2022 Wildfire: Forest zones to west",
        "  \u2022 Drought: Water-cooled system",
        "  \u2022 Water Risk: Coastal cooling intake",
        "",
        "Korean name: \uc0bc\ucca9\ud654\ub825\ubc1c\uc804\uc18c",
    ]
    for d in details:
        p = tf.add_paragraph()
        p.text = d
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(3)


def make_hazard_interpretation(prs):
    """Slide 7: Hazard Interpretation Logic."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Hazard Interpretation Logic")

    add_text_box(slide, "SECTION 2: PHYSICAL RISK METHODOLOGY",
                 Inches(0.5), Inches(1.15), Inches(5), Inches(0.4),
                 font_size=11, bold=True, color=TEAL)

    add_text_box(slide, "How raw PLANiT outputs become cashflow impacts:",
                 Inches(0.5), Inches(1.6), Inches(11), Inches(0.5),
                 font_size=14, color=DARK_GRAY)

    # Conversion table
    hazard_data = [
        ["Hazard", "Raw Output", "Conversion Formula", "Example"],
        ["Wildfire", "AAI (KRW)",
         "outage_rate = AAI / replacement_value (4.879T)",
         "67.9M KRW \u2192 0.001%"],
        ["Drought", "impact_mean",
         "capacity_derate = impact_mean \u00d7 1.0",
         "0.00094 \u2192 0.094%"],
        ["Water Risk", "impact_mean",
         "water_cap = max(0, 1 \u2212 impact_mean)",
         "0.00702 \u2192 99.3% avail."],
    ]
    add_table(slide, hazard_data, [2.0, 2.0, 5.0, 3.0],
              Inches(0.5), Inches(2.2), font_size=12)

    # Compound risk section
    add_text_box(slide, "Compound Risk Multiplier",
                 Inches(0.5), Inches(4.2), Inches(5), Inches(0.5),
                 font_size=18, bold=True, color=NAVY)

    compound_items = [
        "Hazards are not independent \u2014 droughts increase wildfire risk",
        "Compound multiplier range: 1.0x (baseline) to 1.25x (high scenario)",
        "Based on Zscheischler et al. (2018) compound event framework",
        "Total CF Loss = (wildfire + drought + water risk) \u00d7 compound_mult",
    ]
    add_bullet_list(slide, compound_items, Inches(0.5), Inches(4.8),
                    Inches(12), Inches(2.0), font_size=13)

    # Source note
    add_text_box(slide,
                 "Source: src/planit/adapter.py, src/planit/config.py  |  "
                 "Replacement value: KRW 4.879 trillion (KEPCO asset registry)",
                 Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
                 font_size=10, color=RGBColor(0x88, 0x88, 0x88))


def make_physical_scenarios(prs):
    """Slide 8: Physical Risk by SSP Scenario."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Physical Risk by SSP Scenario")

    add_text_box(slide, "SECTION 2: PHYSICAL RISK METHODOLOGY",
                 Inches(0.5), Inches(1.15), Inches(5), Inches(0.4),
                 font_size=11, bold=True, color=TEAL)

    # PLANiT results table
    ssp_data = [
        ["Scenario", "Wildfire", "Drought", "Water Risk", "Total CF Loss"],
        ["Baseline (2024)", "0.001%", "0.094%", "0.000%", "0.095%"],
        ["SSP1-2.6 (2040)", "0.005%", "0.257%", "0.559%", "0.820%"],
        ["SSP5-8.5 (2050)", "0.001%", "0.101%", "0.546%", "0.648%"],
    ]
    add_table(slide, ssp_data, [2.5, 2.0, 2.0, 2.0, 2.0],
              Inches(0.5), Inches(1.8), font_size=13)

    # Comparison: PLANiT vs Literature
    add_text_box(slide, "PLANiT vs Literature Comparison",
                 Inches(0.5), Inches(3.8), Inches(8), Inches(0.5),
                 font_size=18, bold=True, color=NAVY)

    comp_data = [
        ["Source", "Wildfire", "Flood", "Notes"],
        ["PLANiT (site-specific)", "0.001%", "n/a",
         "CLIMADA AAI for Samcheok coordinates"],
        ["Literature (Korea avg)", "0.055%", "0.003%",
         "Kim et al. 2025 / Kang & Lee 2024"],
        ["Ratio", "55x lower", "n/a",
         "PLANiT more granular; site in low-risk zone"],
    ]
    add_table(slide, comp_data, [3.0, 1.5, 1.5, 5.5],
              Inches(0.5), Inches(4.4), font_size=12)

    # Callout
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.2),
        Inches(12.3), Inches(1.0)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
    shape.line.color.rgb = GREEN_ACCENT

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.1)

    p = tf.paragraphs[0]
    p.text = ("Conclusion: Physical risk at Samcheok is modest (< 1% CF loss even under SSP5-8.5). "
              "Transition and policy risks dominate the CRP.")
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1E, 0x8E, 0x4E)


def make_scenario_comparison(prs, scenarios):
    """Slide 9: Scenario Comparison Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Scenario Comparison Overview")

    add_text_box(slide, "SECTION 3: KEY RESULTS",
                 Inches(0.5), Inches(1.15), Inches(5), Inches(0.4),
                 font_size=11, bold=True, color=TEAL)

    # Embed NPV comparison chart
    if FIG1.exists():
        slide.shapes.add_picture(
            str(FIG1), Inches(0.3), Inches(1.7), Inches(7.0), Inches(5.3)
        )

    # Results table on right
    key_scenarios = [
        "baseline", "moderate_transition", "aggressive_transition",
        "combined_aggressive", "low_demand", "enhanced_11th_plan"
    ]
    display_names = {
        "baseline": "Baseline",
        "moderate_transition": "Moderate Transition",
        "aggressive_transition": "Aggressive Transition",
        "combined_aggressive": "Combined Aggressive",
        "low_demand": "Low Demand",
        "enhanced_11th_plan": "Enhanced 11th Plan",
    }

    table_rows = [["Scenario", "NPV (M)", "IRR", "DSCR", "Rating"]]
    for s_name in key_scenarios:
        for s in scenarios:
            if s["scenario"] == s_name:
                npv = float(s["npv_million"])
                irr = float(s["irr_pct"])
                dscr = float(s["avg_dscr"])
                rating = s["overall_rating"]
                table_rows.append([
                    display_names.get(s_name, s_name),
                    f"{npv:,.0f}",
                    f"{irr:.1f}%",
                    f"{dscr:.2f}",
                    rating,
                ])
                break

    add_table(slide, table_rows, [2.8, 1.2, 1.0, 1.0, 1.0],
              Inches(7.5), Inches(1.8), font_size=12)

    # Bottom note
    add_text_box(slide,
                 "NPV in million USD, 8% discount rate. IRR = Project Internal Rate of Return. "
                 "DSCR = Debt Service Coverage Ratio. Rating per KIS methodology.",
                 Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
                 font_size=10, color=RGBColor(0x88, 0x88, 0x88))


def make_rating_migration(prs):
    """Slide 10: Credit Rating Migration."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Credit Rating Migration")

    add_text_box(slide, "SECTION 3: KEY RESULTS",
                 Inches(0.5), Inches(1.15), Inches(5), Inches(0.4),
                 font_size=11, bold=True, color=TEAL)

    # Embed rating migration chart
    if FIG3.exists():
        slide.shapes.add_picture(
            str(FIG3), Inches(0.3), Inches(1.7), Inches(8.0), Inches(5.3)
        )

    # Key migration callouts on right
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.7),
        Inches(4.0), Inches(5.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xEC)
    shape.line.color.rgb = RED_ACCENT

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)

    p = tf.paragraphs[0]
    p.text = "Rating Migration Summary"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY

    migrations = [
        ("\u2191 Baseline: A \u2192 AA", "Upgrade by 1 notch (CRP = -50 bps)", GREEN_ACCENT),
        ("\u2192 Moderate Transition: A \u2192 A", "No change (CRP = 0 bps)", DARK_GRAY),
        ("\u2192 Aggressive Transition: A \u2192 A", "No change (CRP = 0 bps)", DARK_GRAY),
        ("\u2193 Low Demand: A \u2192 BBB", "Downgrade 1 notch (CRP = +85 bps)", RGBColor(0xE6, 0x7E, 0x22)),
        ("\u2193\u2193 Enhanced 11th Plan: A \u2192 CC", "Downgrade 5 notches (CRP = +1,020 bps)", RED_ACCENT),
    ]

    for arrow_text, detail, color in migrations:
        p = tf.add_paragraph()
        p.text = ""
        p.space_before = Pt(10)

        p = tf.add_paragraph()
        p.text = arrow_text
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = color

        p = tf.add_paragraph()
        p.text = detail
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GRAY

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Counterfactual: A (150 bps)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = NAVY


def make_key_findings(prs):
    """Slide 11: Key Findings."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Key Findings & Conclusions")

    add_text_box(slide, "SECTION 3: KEY RESULTS",
                 Inches(0.5), Inches(1.15), Inches(5), Inches(0.4),
                 font_size=11, bold=True, color=TEAL)

    # Finding 1
    shape1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8),
        Inches(6.0), Inches(2.2)
    )
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
    shape1.line.color.rgb = GREEN_ACCENT
    tf = shape1.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = "1. Physical Risk is Modest"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = ("Physical hazards cause only 0.10\u20130.82% capacity factor loss. "
               "Even under SSP5-8.5, the NPV impact is less than USD 62M (2% of baseline NPV). "
               "Physical risk is not the primary driver of climate risk premium.")
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_GRAY
    p2.space_before = Pt(8)

    # Finding 2
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8),
        Inches(6.0), Inches(2.2)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xEC)
    shape2.line.color.rgb = RED_ACCENT
    tf = shape2.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = "2. Transition/Policy Risk Dominates"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = ("Korea's 11th Basic Energy Plan causes CC-rated distress "
               "(NPV \u2212USD 3.0B, DSCR 0.59x). "
               "The aggressive coal phase-out policy alone drives a 5-notch downgrade "
               "and CRP of +1,020 bps.")
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_GRAY
    p2.space_before = Pt(8)

    # Finding 3
    shape3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.3),
        Inches(6.0), Inches(2.2)
    )
    shape3.fill.solid()
    shape3.fill.fore_color.rgb = RGBColor(0xFD, 0xF2, 0xE9)
    shape3.line.color.rgb = RGBColor(0xE6, 0x7E, 0x22)
    tf = shape3.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = "3. CRP Range: \u221250 to +1,020 bps"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = ("The climate risk premium spans from \u221250 bps (baseline upgrade to AA) "
               "to +1,020 bps (enhanced 11th plan, CC distress). "
               "Investment-grade is maintained through combined_aggressive; "
               "breached only under policy acceleration.")
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_GRAY
    p2.space_before = Pt(8)

    # Finding 4
    shape4 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.3),
        Inches(6.0), Inches(2.2)
    )
    shape4.fill.solid()
    shape4.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
    shape4.line.color.rgb = NAVY
    tf = shape4.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = "4. Investment Implications"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = ("Current baseline economics are strong (AA, 12% IRR, DSCR 2.15x). "
               "However, investors must price in transition policy tail risk. "
               "The 11th Basic Plan scenario represents a realistic near-term threat "
               "that could strand the asset entirely.")
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_GRAY
    p2.space_before = Pt(8)

    # Footer
    add_text_box(slide,
                 "Climate Risk Premium Model  |  Samcheok 2,100 MW  |  "
                 "11 scenarios, 40-year DCF  |  February 2026",
                 Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
                 font_size=10, color=RGBColor(0x88, 0x88, 0x88),
                 alignment=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════
# Main
# ═════════════════════════════════════════════════════════════════════════
def main():
    print("Generating CRP Presentation...")

    # Step 1: Generate map
    print("  Generating Samcheok map...")
    generate_map()

    # Step 2: Load data
    print("  Loading data...")
    plant = load_plant_params(PLANT_CSV)
    scenarios = load_csv(SCENARIO_CSV)

    # Step 3: Create presentation (widescreen 16:9)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Build slides
    print("  Building slides...")
    make_title_slide(prs)          # Slide 1
    make_plant_overview(prs, plant) # Slide 2
    make_price_revenue(prs)         # Slide 3
    make_cost_financing(prs)        # Slide 4
    make_planit_framework(prs)      # Slide 5
    make_map_slide(prs)             # Slide 6
    make_hazard_interpretation(prs) # Slide 7
    make_physical_scenarios(prs)    # Slide 8
    make_scenario_comparison(prs, scenarios)  # Slide 9
    make_rating_migration(prs)      # Slide 10
    make_key_findings(prs)          # Slide 11

    # Step 4: Save
    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PPTX))
    size_kb = OUTPUT_PPTX.stat().st_size / 1024
    print(f"\n  Presentation saved: {OUTPUT_PPTX}")
    print(f"  File size: {size_kb:.0f} KB")
    print(f"  Slides: {len(prs.slides)}")
    print("  Done!")


if __name__ == "__main__":
    main()
